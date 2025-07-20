from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create directories if they don't exist
base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
presentation_dir = os.path.join(base_dir, 'reports', 'presentations')
if not os.path.exists(presentation_dir):
    os.makedirs(presentation_dir)

def create_presentation():
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Add title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9), Inches(1))
    tf = title.text_frame
    tf.text = "Marketing Campaign & Product Bundle Optimization"
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.color.rgb = RGBColor(0x4E, 0xCD, 0xC4)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1))
    tf = subtitle.text_frame
    tf.text = "Q4 2023 Performance Review"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add trend plots slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf = title.text_frame
    tf.text = "Performance Trends"
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.color.rgb = RGBColor(0x4E, 0xCD, 0xC4)
    
    # Add trend plots
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12)
    height = Inches(4.5)
    
    try:
        img_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'reports', 'figures', 'presentation', 'trend_plots.png')
        slide.shapes.add_picture(img_path, left, top, width, height)
    except:
        txt = slide.shapes.add_textbox(left, top, width, height)
        tf = txt.text_frame
        tf.text = "Trend Plots (Image not available)"
        tf.paragraphs[0].font.size = Pt(18)
    
    # Add customer map slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf = title.text_frame
    tf.text = "Customer Distribution"
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.color.rgb = RGBColor(0x4E, 0xCD, 0xC4)
    
    # Add customer map
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12)
    height = Inches(4.5)
    
    try:
        img_path = '../reports/figures/presentation/customer_map.png'
        slide.shapes.add_picture(img_path, left, top, width, height)
    except:
        txt = slide.shapes.add_textbox(left, top, width, height)
        tf = txt.text_frame
        tf.text = "Customer Map (Image not available)"
        tf.paragraphs[0].font.size = Pt(18)
    
    # Add segmentation charts slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf = title.text_frame
    tf.text = "Customer Segmentation Analysis"
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.color.rgb = RGBColor(0x4E, 0xCD, 0xC4)
    
    # Add segmentation charts
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12)
    height = Inches(4.5)
    
    try:
        img_path = '../reports/figures/presentation/segmentation_charts.png'
        slide.shapes.add_picture(img_path, left, top, width, height)
    except:
        txt = slide.shapes.add_textbox(left, top, width, height)
        tf = txt.text_frame
        tf.text = "Segmentation Charts (Image not available)"
        tf.paragraphs[0].font.size = Pt(18)
    
    # Add KPIs slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf = title.text_frame
    tf.text = "Key Performance Indicators"
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.color.rgb = RGBColor(0x4E, 0xCD, 0xC4)
    
    # Add KPIs
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12)
    height = Inches(4.5)
    
    try:
        img_path = '../reports/figures/presentation/kpi_bars.png'
        slide.shapes.add_picture(img_path, left, top, width, height)
    except:
        txt = slide.shapes.add_textbox(left, top, width, height)
        tf = txt.text_frame
        tf.text = "KPI Bar Graphs (Image not available)"
        tf.paragraphs[0].font.size = Pt(18)
    
    # Save the presentation
    prs.save(os.path.join(presentation_dir, 'marketing_analysis.pptx'))
    print("Presentation created successfully!")

if __name__ == "__main__":
    create_presentation()
